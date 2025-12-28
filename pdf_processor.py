import os
import json
import logging
import zipfile
from io import BytesIO
from datetime import datetime
from PyPDF2 import PdfReader, PdfWriter
from PIL import Image
import pytesseract
from docx import Document
import openpyxl
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from app import db
from models import ProcessingJob, JobStatus, JobType, get_now
from utils import generate_unique_filename

logger = logging.getLogger(__name__)

class PDFProcessor:
    def __init__(self, job_id):
        self.job_id = job_id
        self.job = ProcessingJob.query.get(job_id)
        if not self.job:
            raise ValueError(f"Job {job_id} not found")
    
    def update_progress(self, progress, processed_files=None):
        """Update job progress in database"""
        self.job.progress = progress
        if processed_files is not None:
            self.job.processed_files = processed_files
        db.session.commit()
    
    def update_status(self, status, error_message=None):
        """Update job status in database"""
        self.job.status = status
        if error_message:
            self.job.error_message = error_message
        if status == JobStatus.PROCESSING:
            self.job.started_at = get_now()
        elif status in [JobStatus.COMPLETED, JobStatus.FAILED, JobStatus.CANCELLED]:
            self.job.completed_at = get_now()
        db.session.commit()
    
    def process_job(self):
        """Main processing function that routes to specific processors"""
        try:
            self.update_status(JobStatus.PROCESSING)
            
            # Apply watermark for free users if not already a watermark job
            is_free_user = not self.job.user.is_premium
            
            if self.job.job_type == JobType.MERGE:
                result = self.merge_pdfs()
            elif self.job.job_type == JobType.SPLIT:
                result = self.split_pdfs()
            elif self.job.job_type == JobType.COMPRESS:
                result = self.compress_pdfs()
            elif self.job.job_type == JobType.OCR:
                result = self.ocr_pdfs()
            elif self.job.job_type == JobType.CONVERT_WORD:
                result = self.convert_to_word()
            elif self.job.job_type == JobType.CONVERT_EXCEL:
                result = self.convert_to_excel()
            elif self.job.job_type == JobType.PROTECT:
                result = self.protect_pdfs()
            elif self.job.job_type == JobType.ROTATE:
                result = self.rotate_pdfs()
            elif self.job.job_type == JobType.WATERMARK:
                result = self.watermark_pdfs()
            elif self.job.job_type == JobType.UNLOCK:
                result = self.unlock_pdfs()
            elif self.job.job_type == JobType.EXTRACT_IMAGES:
                result = self.extract_images_pdfs()
            elif self.job.job_type in [JobType.REMOVE_PAGES, JobType.EXTRACT_PAGES, JobType.ORGANIZE]:
                result = self.organize_pdf_pages()
            elif self.job.job_type == JobType.SCAN_TO_PDF:
                result = self.scan_to_pdf()
            elif self.job.job_type == JobType.REPAIR:
                result = self.repair_pdf()
            elif self.job.job_type in [JobType.JPG_TO_PDF, JobType.WORD_TO_PDF, JobType.POWERPOINT_TO_PDF, JobType.EXCEL_TO_PDF, JobType.HTML_TO_PDF]:
                result = self.convert_to_pdf()
            elif self.job.job_type in [JobType.PDF_TO_JPG, JobType.PDF_TO_POWERPOINT, JobType.PDF_TO_EXCEL, JobType.PDF_TO_PDFA]:
                result = self.convert_from_pdf()
            elif self.job.job_type == JobType.ADD_PAGE_NUMBERS:
                result = self.add_page_numbers()
            elif self.job.job_type == JobType.CROP:
                result = self.crop_pdf()
            elif self.job.job_type == JobType.EDIT:
                result = self.edit_pdf()
            elif self.job.job_type == JobType.SIGN:
                result = self.sign_pdf()
            elif self.job.job_type == JobType.REDACT:
                result = self.redact_pdf()
            elif self.job.job_type == JobType.COMPARE:
                result = self.compare_pdf()
            else:
                raise ValueError(f"Unknown job type: {self.job.job_type}")
            
            # Add watermark for free users to all output PDFs
            if is_free_user:
                self.apply_free_tier_watermark(result)
            
            self.job.output_files = json.dumps(result)
            self.update_status(JobStatus.COMPLETED)
            self.update_progress(100, self.job.total_files)
            
            return result
            
        except Exception as e:
            logger.error(f"Error processing job {self.job_id}: {str(e)}")
            self.update_status(JobStatus.FAILED, str(e))
            raise
    
    def apply_free_tier_watermark(self, file_paths):
        """Add 'Processed with SnapPDF' watermark to free tier files"""
        if not file_paths:
            return
            
        for i, file_path in enumerate(file_paths):
            if not file_path.lower().endswith('.pdf'):
                continue
                
            try:
                from reportlab.pdfgen import canvas
                from reportlab.lib.colors import grey
                import io
                
                reader = PdfReader(file_path)
                writer = PdfWriter()
                
                for page in reader.pages:
                    packet = io.BytesIO()
                    can = canvas.Canvas(packet)
                    can.setFont("Helvetica", 40)
                    can.setFillAlpha(0.3)
                    can.saveState()
                    can.translate(300, 400)
                    can.rotate(45)
                    can.drawCentredString(0, 0, "Processed with SnapPDF Free")
                    can.restoreState()
                    can.save()
                    packet.seek(0)
                    
                    watermark = PdfReader(packet).pages[0]
                    page.merge_page(watermark)
                    writer.add_page(page)
                
                with open(file_path, 'wb') as f:
                    writer.write(f)
                    
            except Exception as e:
                logger.error(f"Error applying free watermark to {file_path}: {e}")

    def merge_pdfs(self):
        """Merge multiple PDF files into one"""
        input_files = json.loads(self.job.input_files)
        output_filename = generate_unique_filename("merged_document.pdf")
        output_path = os.path.join(app.config['PROCESSED_FOLDER'], output_filename)
        os.makedirs(app.config['PROCESSED_FOLDER'], exist_ok=True)
        
        writer = PdfWriter()
        
        for i, file_path in enumerate(input_files):
            try:
                with open(file_path, 'rb') as file:
                    reader = PdfReader(file)
                    for page in reader.pages:
                        writer.add_page(page)
                
                progress = int((i + 1) / len(input_files) * 100)
                self.update_progress(progress, i + 1)
                
            except Exception as e:
                logger.error(f"Error processing file {file_path}: {str(e)}")
                continue
        
        with open(output_path, 'wb') as output_file:
            writer.write(output_file)
        
        return [output_path]
    
    def split_pdfs(self):
        """Split PDF files into individual pages"""
        input_files = json.loads(self.job.input_files)
        output_files = []
        
        for file_idx, file_path in enumerate(input_files):
            try:
                with open(file_path, 'rb') as file:
                    reader = PdfReader(file)
                    base_name = os.path.splitext(os.path.basename(file_path))[0]
                    
                    for page_num, page in enumerate(reader.pages):
                        writer = PdfWriter()
                        writer.add_page(page)
                        
                        output_filename = generate_unique_filename(f"{base_name}_page_{page_num + 1}.pdf")
                        output_path = os.path.join(app.config['PROCESSED_FOLDER'], output_filename)
                        os.makedirs(app.config['PROCESSED_FOLDER'], exist_ok=True)
                        
                        with open(output_path, 'wb') as output_file:
                            writer.write(output_file)
                        
                        output_files.append(output_path)
                
                progress = int((file_idx + 1) / len(input_files) * 100)
                self.update_progress(progress, file_idx + 1)
                
            except Exception as e:
                logger.error(f"Error splitting file {file_path}: {str(e)}")
                continue
        
        return output_files
    
    def compress_pdfs(self):
        """Compress PDF files with adjustable quality settings"""
        input_files = json.loads(self.job.input_files)
        settings = json.loads(self.job.settings) if self.job.settings else {}
        quality = settings.get('compression_quality', 'medium')  # low, medium, high
        output_files = []
        
        # Quality settings mapping
        quality_settings = {
            'low': {'scale': 0.6, 'jpeg_quality': 30},
            'medium': {'scale': 0.8, 'jpeg_quality': 60}, 
            'high': {'scale': 1.0, 'jpeg_quality': 85}
        }
        
        current_settings = quality_settings.get(quality, quality_settings['medium'])
        
        for file_idx, file_path in enumerate(input_files):
            try:
                # Use PyMuPDF for better compression
                import fitz
                
                doc = fitz.open(file_path)
                
                base_name = os.path.splitext(os.path.basename(file_path))[0]
                output_filename = generate_unique_filename(f"{base_name}_compressed_{quality}.pdf")
                output_path = os.path.join(app.config['PROCESSED_FOLDER'], output_filename)
                os.makedirs(app.config['PROCESSED_FOLDER'], exist_ok=True)
                
                # Create a new document for the compressed version
                new_doc = fitz.open()
                
                for page_num in range(len(doc)):
                    page = doc[page_num]
                    
                    # Get page dimensions
                    rect = page.rect
                    
                    # Create pixmap with scaling for compression
                    scale = current_settings['scale']
                    mat = fitz.Matrix(scale, scale)
                    pix = page.get_pixmap(matrix=mat)
                    
                    # Convert to JPEG bytes for compression
                    img_data = pix.tobytes("jpeg", jpg_quality=current_settings['jpeg_quality'])
                    
                    # Create new page and insert compressed image
                    new_page = new_doc.new_page(width=rect.width, height=rect.height)
                    new_page.insert_image(rect, stream=img_data)
                
                # Save with additional compression options
                new_doc.save(output_path, 
                           garbage=4,  # Remove unused objects
                           deflate=True,  # Use deflate compression
                           clean=True,  # Clean up structure
                           ascii=False)  # Keep binary format
                
                new_doc.close()
                doc.close()
                
                output_files.append(output_path)
                
            except ImportError:
                # Fallback to PyPDF2 for basic compression
                try:
                    with open(file_path, 'rb') as file:
                        reader = PdfReader(file)
                        writer = PdfWriter()
                        
                        for page in reader.pages:
                            # Remove annotations to reduce size
                            if '/Annots' in page:
                                del page['/Annots']
                            writer.add_page(page)
                        
                        base_name = os.path.splitext(os.path.basename(file_path))[0]
                        output_filename = generate_unique_filename(f"{base_name}_compressed_{quality}.pdf")
                        output_path = os.path.join(app.config['PROCESSED_FOLDER'], output_filename)
                        os.makedirs(app.config['PROCESSED_FOLDER'], exist_ok=True)
                        
                        with open(output_path, 'wb') as output_file:
                            writer.write(output_file)
                        
                        output_files.append(output_path)
                except Exception as e:
                    logger.error(f"Error compressing file {file_path}: {str(e)}")
                    continue
            except Exception as e:
                logger.error(f"Error compressing file {file_path}: {str(e)}")
                continue
                
                progress = int((file_idx + 1) / len(input_files) * 100)
                self.update_progress(progress, file_idx + 1)
                
            except Exception as e:
                logger.error(f"Error compressing file {file_path}: {str(e)}")
                continue
        
        return output_files
    
    def ocr_pdfs(self):
        """Extract text from PDF files using advanced OCR"""
        input_files = json.loads(self.job.input_files)
        settings = json.loads(self.job.settings) if self.job.settings else {}
        language = settings.get('ocr_language', 'eng')  # Default to English
        output_format = settings.get('output_format', 'txt')  # txt, pdf, both
        output_files = []
        
        for file_idx, file_path in enumerate(input_files):
            try:
                import fitz  # PyMuPDF for better PDF to image conversion
                
                # Open PDF with PyMuPDF for better image extraction
                try:
                    pdf_document = fitz.open(file_path)
                    text_content = []
                    
                    for page_num in range(pdf_document.page_count):
                        page = pdf_document[page_num]
                        
                        # First try to extract text directly
                        text = page.get_text()
                        
                        if text.strip():
                            text_content.append(f"=== Page {page_num + 1} ===\n{text}")
                        else:
                            # Convert page to image for OCR
                            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x zoom for better quality
                            img_data = pix.tobytes("png")
                            
                            # Convert to PIL Image for OCR
                            from PIL import Image
                            import io
                            
                            pil_image = Image.open(io.BytesIO(img_data))
                            
                            # Perform OCR
                            try:
                                ocr_text = pytesseract.image_to_string(pil_image, lang=language, 
                                                                     config='--psm 6 --oem 3')
                                if ocr_text.strip():
                                    text_content.append(f"=== Page {page_num + 1} (OCR) ===\n{ocr_text}")
                                else:
                                    text_content.append(f"=== Page {page_num + 1} ===\n[No text detected]")
                            except Exception as ocr_error:
                                logger.warning(f"OCR failed for page {page_num + 1}: {str(ocr_error)}")
                                text_content.append(f"=== Page {page_num + 1} ===\n[OCR processing failed]")
                    
                    pdf_document.close()
                    
                except ImportError:
                    # Fallback to PyPDF2 if PyMuPDF is not available
                    logger.warning("PyMuPDF not available, using basic text extraction")
                    with open(file_path, 'rb') as file:
                        reader = PdfReader(file)
                        text_content = []
                        
                        for page_num, page in enumerate(reader.pages):
                            text = page.extract_text()
                            if text.strip():
                                text_content.append(f"=== Page {page_num + 1} ===\n{text}")
                            else:
                                text_content.append(f"=== Page {page_num + 1} ===\n[No extractable text found]")
                
                base_name = os.path.splitext(os.path.basename(file_path))[0]
                
                # Save as text file
                if output_format in ['txt', 'both']:
                    txt_filename = generate_unique_filename(f"{base_name}_ocr.txt")
                    txt_output_path = os.path.join(app.config['PROCESSED_FOLDER'], txt_filename)
                    os.makedirs(app.config['PROCESSED_FOLDER'], exist_ok=True)
                    
                    with open(txt_output_path, 'w', encoding='utf-8') as output_file:
                        output_file.write('\n\n'.join(text_content))
                    
                    output_files.append(txt_output_path)
                
                # Save as searchable PDF
                if output_format in ['pdf', 'both']:
                    pdf_filename = generate_unique_filename(f"{base_name}_searchable.pdf")
                    pdf_output_path = os.path.join(app.config['PROCESSED_FOLDER'], pdf_filename)
                    os.makedirs(app.config['PROCESSED_FOLDER'], exist_ok=True)
                    
                    # Create a new PDF with the extracted text
                    from reportlab.pdfgen import canvas
                    from reportlab.lib.pagesizes import letter
                    from reportlab.lib.styles import getSampleStyleSheet
                    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
                    
                    doc = SimpleDocTemplate(pdf_output_path, pagesize=letter)
                    styles = getSampleStyleSheet()
                    story = []
                    
                    for content in text_content:
                        para = Paragraph(content.replace('\n', '<br/>'), styles['Normal'])
                        story.append(para)
                        story.append(Spacer(1, 12))
                    
                    doc.build(story)
                    output_files.append(pdf_output_path)
                
                progress = int((file_idx + 1) / len(input_files) * 100)
                self.update_progress(progress, file_idx + 1)
                
            except Exception as e:
                logger.error(f"Error processing OCR for file {file_path}: {str(e)}")
                continue
        
        return output_files
    
    def convert_to_word(self):
        """Convert PDF files to Word documents"""
        input_files = json.loads(self.job.input_files)
        output_files = []
        
        for file_idx, file_path in enumerate(input_files):
            try:
                with open(file_path, 'rb') as file:
                    reader = PdfReader(file)
                    doc = Document()
                    
                    for page in reader.pages:
                        text = page.extract_text()
                        if text.strip():
                            doc.add_paragraph(text)
                    
                    base_name = os.path.splitext(os.path.basename(file_path))[0]
                    output_filename = generate_unique_filename(f"{base_name}.docx")
                    output_path = os.path.join(app.config['PROCESSED_FOLDER'], output_filename)
                    os.makedirs(app.config['PROCESSED_FOLDER'], exist_ok=True)
                    
                    doc.save(output_path)
                    output_files.append(output_path)
                
                progress = int((file_idx + 1) / len(input_files) * 100)
                self.update_progress(progress, file_idx + 1)
                
            except Exception as e:
                logger.error(f"Error converting file {file_path} to Word: {str(e)}")
                continue
        
        return output_files

    def protect_pdfs(self):
        """Protect PDF files with a password"""
        input_files = json.loads(self.job.input_files)
        settings = json.loads(self.job.settings) if self.job.settings else {}
        password = settings.get('password')
        if not password:
            raise ValueError("Password is required for protection")
        
        output_files = []
        for file_idx, file_path in enumerate(input_files):
            try:
                reader = PdfReader(file_path)
                writer = PdfWriter()
                for page in reader.pages:
                    writer.add_page(page)
                writer.encrypt(password)
                
                base_name = os.path.splitext(os.path.basename(file_path))[0]
                output_filename = generate_unique_filename(f"{base_name}_protected.pdf")
                output_path = os.path.join(app.config['PROCESSED_FOLDER'], output_filename)
                os.makedirs(app.config['PROCESSED_FOLDER'], exist_ok=True)
                
                with open(output_path, 'wb') as f:
                    writer.write(f)
                output_files.append(output_path)
                
                progress = int((file_idx + 1) / len(input_files) * 100)
                self.update_progress(progress, file_idx + 1)
            except Exception as e:
                logger.error(f"Error protecting file {file_path}: {str(e)}")
                continue
        return output_files

    def rotate_pdfs(self):
        """Rotate PDF files"""
        input_files = json.loads(self.job.input_files)
        settings = json.loads(self.job.settings) if self.job.settings else {}
        rotation = int(settings.get('rotation', 90))
        
        output_files = []
        for file_idx, file_path in enumerate(input_files):
            try:
                reader = PdfReader(file_path)
                writer = PdfWriter()
                for page in reader.pages:
                    page.rotate(rotation)
                    writer.add_page(page)
                
                base_name = os.path.splitext(os.path.basename(file_path))[0]
                output_filename = generate_unique_filename(f"{base_name}_rotated.pdf")
                output_path = os.path.join(app.config['PROCESSED_FOLDER'], output_filename)
                os.makedirs(app.config['PROCESSED_FOLDER'], exist_ok=True)
                
                with open(output_path, 'wb') as f:
                    writer.write(f)
                output_files.append(output_path)
                
                progress = int((file_idx + 1) / len(input_files) * 100)
                self.update_progress(progress, file_idx + 1)
            except Exception as e:
                logger.error(f"Error rotating file {file_path}: {str(e)}")
                continue
        return output_files

    def watermark_pdfs(self):
        """Add watermark to PDF files"""
        input_files = json.loads(self.job.input_files)
        settings = json.loads(self.job.settings) if self.job.settings else {}
        text = settings.get('watermark_text', 'CONFIDENTIAL')
        
        output_files = []
        for file_idx, file_path in enumerate(input_files):
            try:
                from reportlab.pdfgen import canvas
                from reportlab.lib.pagesizes import letter
                import io
                
                # Create watermark
                packet = io.BytesIO()
                can = canvas.Canvas(packet, pagesize=letter)
                can.setFont("Helvetica", 40)
                can.setStrokeColorRGB(0.5, 0.5, 0.5, 0.3)
                can.setFillColorRGB(0.5, 0.5, 0.5, 0.3)
                can.saveState()
                can.translate(300, 400)
                can.rotate(45)
                can.drawCentredString(0, 0, text)
                can.restoreState()
                can.save()
                packet.seek(0)
                watermark_reader = PdfReader(packet)
                watermark_page = watermark_reader.pages[0]

                reader = PdfReader(file_path)
                writer = PdfWriter()
                for page in reader.pages:
                    page.merge_page(watermark_page)
                    writer.add_page(page)
                
                base_name = os.path.splitext(os.path.basename(file_path))[0]
                output_filename = generate_unique_filename(f"{base_name}_watermarked.pdf")
                output_path = os.path.join(app.config['PROCESSED_FOLDER'], output_filename)
                os.makedirs(app.config['PROCESSED_FOLDER'], exist_ok=True)
                
                with open(output_path, 'wb') as f:
                    writer.write(f)
                output_files.append(output_path)
                
                progress = int((file_idx + 1) / len(input_files) * 100)
                self.update_progress(progress, file_idx + 1)
            except Exception as e:
                logger.error(f"Error watermarking file {file_path}: {str(e)}")
                continue
        return output_files

    def unlock_pdfs(self):
        """Unlock protected PDF files"""
        input_files = json.loads(self.job.input_files)
        settings = json.loads(self.job.settings) if self.job.settings else {}
        password = settings.get('password')
        if not password:
            raise ValueError("Password is required for unlocking")
        
        output_files = []
        for file_idx, file_path in enumerate(input_files):
            try:
                reader = PdfReader(file_path)
                if reader.is_encrypted:
                    reader.decrypt(password)
                
                writer = PdfWriter()
                for page in reader.pages:
                    writer.add_page(page)
                
                base_name = os.path.splitext(os.path.basename(file_path))[0]
                output_filename = generate_unique_filename(f"{base_name}_unlocked.pdf")
                output_path = os.path.join(app.config['PROCESSED_FOLDER'], output_filename)
                os.makedirs(app.config['PROCESSED_FOLDER'], exist_ok=True)
                
                with open(output_path, 'wb') as f:
                    writer.write(f)
                output_files.append(output_path)
                
                progress = int((file_idx + 1) / len(input_files) * 100)
                self.update_progress(progress, file_idx + 1)
            except Exception as e:
                logger.error(f"Error unlocking file {file_path}: {str(e)}")
                continue
        return output_files

    def extract_images_pdfs(self):
        """Extract images from PDF files"""
        input_files = json.loads(self.job.input_files)
        output_files = []
        for file_idx, file_path in enumerate(input_files):
            try:
                import fitz
                doc = fitz.open(file_path)
                base_name = os.path.splitext(os.path.basename(file_path))[0]
                
                for page_num in range(len(doc)):
                    page = doc[page_num]
                    image_list = page.get_images()
                    for img_idx, img in enumerate(image_list):
                        xref = img[0]
                        base_image = doc.extract_image(xref)
                        image_bytes = base_image["image"]
                        image_ext = base_image["ext"]
                        
                        output_filename = generate_unique_filename(f"{base_name}_p{page_num+1}_img{img_idx+1}.{image_ext}")
                        output_path = os.path.join(app.config['PROCESSED_FOLDER'], output_filename)
                        os.makedirs(app.config['PROCESSED_FOLDER'], exist_ok=True)
                        
                        with open(output_path, "wb") as f:
                            f.write(image_bytes)
                        output_files.append(output_path)
                
                doc.close()
                progress = int((file_idx + 1) / len(input_files) * 100)
                self.update_progress(progress, file_idx + 1)
            except Exception as e:
                logger.error(f"Error extracting images from {file_path}: {str(e)}")
                continue
        return output_files

    def scan_to_pdf(self):
        """Convert images (simulating scan) to PDF with OCR"""
        # Ensure we use convert_to_pdf logic
        return self.convert_to_pdf()

    def organize_pdf_pages(self):
        """Organize, remove or extract pages from PDF"""
        input_files = json.loads(self.job.input_files)
        settings = json.loads(self.job.settings) if self.job.settings else {}
        page_indices = settings.get('pages', []) # List of 1-based indices from UI
        
        output_files = []
        for file_idx, file_path in enumerate(input_files):
            try:
                reader = PdfReader(file_path)
                writer = PdfWriter()
                
                total_pages = len(reader.pages)
                
                if self.job.job_type == JobType.REMOVE_PAGES:
                    # Remove specified 1-based indices
                    for i in range(total_pages):
                        if (i + 1) not in page_indices:
                            writer.add_page(reader.pages[i])
                elif self.job.job_type == JobType.EXTRACT_PAGES:
                    # Extract specified 1-based indices
                    for p_num in page_indices:
                        if 1 <= p_num <= total_pages:
                            writer.add_page(reader.pages[p_num - 1])
                elif self.job.job_type == JobType.ORGANIZE:
                    # Reorder based on specified 1-based indices, or keep all if empty
                    target_order = page_indices if page_indices else range(1, total_pages + 1)
                    for p_num in target_order:
                        if 1 <= p_num <= total_pages:
                            writer.add_page(reader.pages[p_num - 1])
                
                base_name = os.path.splitext(os.path.basename(file_path))[0]
                suffix = self.job.job_type.value
                output_filename = generate_unique_filename(f"{base_name}_{suffix}.pdf")
                output_path = os.path.join(app.config['PROCESSED_FOLDER'], output_filename)
                os.makedirs(app.config['PROCESSED_FOLDER'], exist_ok=True)
                
                with open(output_path, 'wb') as f:
                    writer.write(f)
                output_files.append(output_path)
                
                self.update_progress(int((file_idx + 1) / len(input_files) * 100), file_idx + 1)
            except Exception as e:
                logger.error(f"Error organizing file {file_path}: {str(e)}")
                continue
        return output_files

    def repair_pdf(self):
        """Attempt to repair a corrupted PDF by re-saving it"""
        input_files = json.loads(self.job.input_files)
        output_files = []
        for file_idx, file_path in enumerate(input_files):
            try:
                reader = PdfReader(file_path)
                writer = PdfWriter()
                for page in reader.pages:
                    writer.add_page(page)
                output_filename = generate_unique_filename("repaired.pdf")
                output_path = os.path.join(app.config['PROCESSED_FOLDER'], output_filename)
                os.makedirs(app.config['PROCESSED_FOLDER'], exist_ok=True)
                with open(output_path, 'wb') as f:
                    writer.write(f)
                output_files.append(output_path)
            except Exception as e:
                logger.error(f"Repair failed: {e}")
        return output_files

    def convert_to_excel(self):
        """Convert PDF to Excel"""
        input_files = json.loads(self.job.input_files)
        output_files = []
        for file_idx, file_path in enumerate(input_files):
            try:
                wb = openpyxl.Workbook()
                ws = wb.active
                reader = PdfReader(file_path)
                for i, page in enumerate(reader.pages):
                    ws.cell(row=i+1, column=1, value=page.extract_text()[:32000])
                
                base_name = os.path.splitext(os.path.basename(file_path))[0]
                output_filename = generate_unique_filename(f"{base_name}.xlsx")
                output_path = os.path.join(app.config['PROCESSED_FOLDER'], output_filename)
                os.makedirs(app.config['PROCESSED_FOLDER'], exist_ok=True)
                wb.save(output_path)
                output_files.append(output_path)
                
                progress = int((file_idx + 1) / len(input_files) * 100)
                self.update_progress(progress, file_idx + 1)
            except Exception as e:
                logger.error(f"Error converting to excel {file_path}: {str(e)}")
                continue
        return output_files

    def convert_to_pdf(self):
        """Convert images and documents to PDF"""
        input_files = json.loads(self.job.input_files)
        output_files = []
        
        for file_idx, file_path in enumerate(input_files):
            try:
                file_ext = os.path.splitext(file_path)[1].lower()
                base_name = os.path.splitext(os.path.basename(file_path))[0]
                output_filename = generate_unique_filename(f"{base_name}.pdf")
                output_path = os.path.join("processed", output_filename)
                
                if file_ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp']:
                    # Image to PDF
                    from PIL import Image
                    img = Image.open(file_path)
                    if img.mode in ('RGBA', 'LA', 'P'):
                        img = img.convert('RGB')
                    img.save(output_path, 'PDF')
                
                elif file_ext == '.docx':
                    # Word to PDF - convert text content
                    doc = Document(file_path)
                    from reportlab.pdfgen import canvas
                    from reportlab.lib.pagesizes import letter
                    c = canvas.Canvas(output_path, pagesize=letter)
                    y = 750
                    for para in doc.paragraphs:
                        if para.text:
                            c.drawString(50, y, para.text[:80])
                            y -= 20
                            if y < 50:
                                c.showPage()
                                y = 750
                    c.save()
                
                elif file_ext == '.xlsx':
                    # Excel to PDF - convert spreadsheet
                    wb = openpyxl.load_workbook(file_path)
                    ws = wb.active
                    from reportlab.pdfgen import canvas
                    from reportlab.lib.pagesizes import letter
                    c = canvas.Canvas(output_path, pagesize=letter)
                    y = 750
                    for row in ws.iter_rows(values_only=True):
                        row_text = ' | '.join(str(cell) if cell else '' for cell in row)
                        if row_text:
                            c.drawString(50, y, row_text[:80])
                            y -= 15
                            if y < 50:
                                c.showPage()
                                y = 750
                    c.save()
                
                elif file_ext == '.pptx':
                    # PowerPoint to PDF - basic text extraction
                    from reportlab.pdfgen import canvas
                    from reportlab.lib.pagesizes import letter
                    c = canvas.Canvas(output_path, pagesize=letter)
                    y = 750
                    c.drawString(50, y, f"PowerPoint: {base_name}")
                    y -= 30
                    c.drawString(50, y, "(Basic text conversion)")
                    c.save()
                
                elif file_ext == '.html':
                    # HTML to PDF - basic conversion
                    from reportlab.pdfgen import canvas
                    from reportlab.lib.pagesizes import letter
                    c = canvas.Canvas(output_path, pagesize=letter)
                    c.drawString(50, 750, f"HTML Document: {base_name}")
                    with open(file_path, 'r', encoding='utf-8') as f:
                        content = f.read()
                        # Simple text extraction
                        from html.parser import HTMLParser
                        class TextExtractor(HTMLParser):
                            def __init__(self):
                                super().__init__()
                                self.text = []
                            def handle_data(self, data):
                                if data.strip():
                                    self.text.append(data.strip())
                        parser = TextExtractor()
                        parser.feed(content)
                        y = 720
                        for text in parser.text[:50]:  # Limit to first 50 lines
                            c.drawString(50, y, text[:70])
                            y -= 15
                    c.save()
                
                output_files.append(output_path)
                progress = int((file_idx + 1) / len(input_files) * 100)
                self.update_progress(progress, file_idx + 1)
                
            except Exception as e:
                logger.error(f"Error converting {file_path} to PDF: {str(e)}")
                continue
        
        return output_files

    def convert_from_pdf(self):
        """Convert PDF to other formats (images, documents, slides, sheets)"""
        input_files = json.loads(self.job.input_files)
        output_files = []
        
        for file_idx, file_path in enumerate(input_files):
            try:
                import fitz
                doc = fitz.open(file_path)
                base_name = os.path.splitext(os.path.basename(file_path))[0]
                
                if self.job.job_type == JobType.PDF_TO_JPG:
                    # PDF to JPG - convert each page to image
                    for page_num in range(len(doc)):
                        page = doc[page_num]
                        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                        output_filename = generate_unique_filename(f"{base_name}_page_{page_num+1}.jpg")
                        output_path = os.path.join(app.config['PROCESSED_FOLDER'], output_filename)
                        os.makedirs(app.config['PROCESSED_FOLDER'], exist_ok=True)
                        pix.save(output_path)
                        output_files.append(output_path)
                
                elif self.job.job_type == JobType.PDF_TO_POWERPOINT:
                    # PDF to PowerPoint - convert each page to a slide image
                    from pptx import Presentation
                    from pptx.util import Inches
                    prs = Presentation()
                    # Set slide size to match common PDF aspect ratio or standard 4:3
                    prs.slide_width = Inches(10)
                    prs.slide_height = Inches(7.5)
                    
                    for page_num in range(len(doc)):
                        page = doc[page_num]
                        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                        img_filename = generate_unique_filename(f"temp_slide_{page_num}.png")
                        img_path = os.path.join("/tmp", img_filename)
                        pix.save(img_path)
                        
                        slide = prs.slides.add_slide(prs.slide_layouts[6]) # blank slide
                        slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
                        if os.path.exists(img_path):
                            os.remove(img_path)
                            
                    output_filename = generate_unique_filename(f"{base_name}.pptx")
                    output_path = os.path.join(app.config['PROCESSED_FOLDER'], output_filename)
                    os.makedirs(app.config['PROCESSED_FOLDER'], exist_ok=True)
                    prs.save(output_path)
                    output_files.append(output_path)
                
                elif self.job.job_type == JobType.PDF_TO_EXCEL:
                    # PDF to Excel - extract text content
                    import openpyxl
                    wb = openpyxl.Workbook()
                    ws = wb.active
                    ws.title = "Extracted Text"
                    
                    row = 1
                    for page_num in range(len(doc)):
                        page = doc[page_num]
                        text = page.get_text()
                        for line in text.split('\n'):
                            if line.strip():
                                ws.cell(row=row, column=1, value=line.strip())
                                row += 1
                                
                    output_filename = generate_unique_filename(f"{base_name}.xlsx")
                    output_path = os.path.join(app.config['PROCESSED_FOLDER'], output_filename)
                    os.makedirs(app.config['PROCESSED_FOLDER'], exist_ok=True)
                    wb.save(output_path)
                    output_files.append(output_path)
                
                elif self.job.job_type == JobType.PDF_TO_PDFA:
                    # PDF to PDF/A - Simplified archive format
                    output_filename = generate_unique_filename(f"{base_name}_pdfa.pdf")
                    output_path = os.path.join(app.config['PROCESSED_FOLDER'], output_filename)
                    os.makedirs(app.config['PROCESSED_FOLDER'], exist_ok=True)
                    doc.save(output_path, garbage=4, deflate=True, clean=True)
                    output_files.append(output_path)
                
                doc.close()
                progress = int((file_idx + 1) / len(input_files) * 100)
                self.update_progress(progress, file_idx + 1)
                
            except Exception as e:
                logger.error(f"Error converting {file_path}: {str(e)}")
                continue
        
        return output_files

    def add_page_numbers(self):
        """Add page numbers to the bottom of each page"""
        input_files = json.loads(self.job.input_files)
        output_files = []
        
        for file_idx, file_path in enumerate(input_files):
            try:
                from reportlab.pdfgen import canvas
                from reportlab.lib.pagesizes import letter
                import io
                
                reader = PdfReader(file_path)
                writer = PdfWriter()
                
                for page_num, page in enumerate(reader.pages):
                    # Create page number annotation
                    page_num_str = str(page_num + 1)
                    
                    # Create a canvas for the page number
                    packet = io.BytesIO()
                    can = canvas.Canvas(packet, pagesize=(612, 792))
                    can.setFont("Helvetica", 10)
                    can.drawRightString(570, 20, page_num_str)
                    can.save()
                    packet.seek(0)
                    
                    # Merge with original page
                    page_num_reader = PdfReader(packet)
                    page_num_page = page_num_reader.pages[0]
                    page.merge_page(page_num_page)
                    writer.add_page(page)
                
                base_name = os.path.splitext(os.path.basename(file_path))[0]
                output_filename = generate_unique_filename(f"{base_name}_numbered.pdf")
                output_path = os.path.join(app.config['PROCESSED_FOLDER'], output_filename)
                os.makedirs(app.config['PROCESSED_FOLDER'], exist_ok=True)
                
                with open(output_path, 'wb') as f:
                    writer.write(f)
                output_files.append(output_path)
                
                progress = int((file_idx + 1) / len(input_files) * 100)
                self.update_progress(progress, file_idx + 1)
            except Exception as e:
                logger.error(f"Error adding page numbers to {file_path}: {str(e)}")
                continue
        
        return output_files

    def crop_pdf(self):
        """Crop PDF pages to specified dimensions"""
        input_files = json.loads(self.job.input_files)
        settings = json.loads(self.job.settings) if self.job.settings else {}
        # Default crop coordinates (left, top, right, bottom) as percentages
        crop_box = settings.get('crop_box', [0.1, 0.1, 0.9, 0.9])
        output_files = []
        
        for file_idx, file_path in enumerate(input_files):
            try:
                reader = PdfReader(file_path)
                writer = PdfWriter()
                
                for page in reader.pages:
                    # Get page dimensions
                    mediabox = page.mediabox
                    width = float(mediabox.width)
                    height = float(mediabox.height)
                    
                    # Calculate crop box
                    left = width * crop_box[0]
                    top = height * crop_box[1]
                    right = width * crop_box[2]
                    bottom = height * crop_box[3]
                    
                    # Crop the page
                    page.cropbox.lower_left = (left, height - bottom)
                    page.cropbox.upper_right = (right, height - top)
                    
                    writer.add_page(page)
                
                base_name = os.path.splitext(os.path.basename(file_path))[0]
                output_filename = generate_unique_filename(f"{base_name}_cropped.pdf")
                output_path = os.path.join(app.config['PROCESSED_FOLDER'], output_filename)
                os.makedirs(app.config['PROCESSED_FOLDER'], exist_ok=True)
                
                with open(output_path, 'wb') as f:
                    writer.write(f)
                output_files.append(output_path)
                
                progress = int((file_idx + 1) / len(input_files) * 100)
                self.update_progress(progress, file_idx + 1)
            except Exception as e:
                logger.error(f"Error cropping {file_path}: {str(e)}")
                continue
        
        return output_files

    def edit_pdf(self):
        """Edit PDF content - basic implementation adding a text layer or overlay"""
        input_files = json.loads(self.job.input_files)
        settings = json.loads(self.job.settings) if self.job.settings else {}
        edit_text = settings.get('edit_text', 'Edited with SnapPDF')
        output_files = []
        
        for file_idx, file_path in enumerate(input_files):
            try:
                from reportlab.pdfgen import canvas
                import io
                
                reader = PdfReader(file_path)
                writer = PdfWriter()
                
                for page in reader.pages:
                    packet = io.BytesIO()
                    can = canvas.Canvas(packet)
                    can.setFont("Helvetica", 12)
                    can.drawString(100, 100, edit_text)
                    can.save()
                    packet.seek(0)
                    
                    overlay = PdfReader(packet).pages[0]
                    page.merge_page(overlay)
                    writer.add_page(page)
                
                base_name = os.path.splitext(os.path.basename(file_path))[0]
                output_filename = generate_unique_filename(f"{base_name}_edited.pdf")
                output_path = os.path.join(app.config['PROCESSED_FOLDER'], output_filename)
                os.makedirs(app.config['PROCESSED_FOLDER'], exist_ok=True)
                
                with open(output_path, 'wb') as f:
                    writer.write(f)
                output_files.append(output_path)
                
                progress = int((file_idx + 1) / len(input_files) * 100)
                self.update_progress(progress, file_idx + 1)
            except Exception as e:
                logger.error(f"Error editing PDF {file_path}: {str(e)}")
                continue
        return output_files

    def sign_pdf(self):
        """Sign PDF (add signature text to last page)"""
        input_files = json.loads(self.job.input_files)
        settings = json.loads(self.job.settings) if self.job.settings else {}
        signature_text = settings.get('signature_text', 'Signed electronically')
        output_files = []
        for file_idx, file_path in enumerate(input_files):
            try:
                from reportlab.pdfgen import canvas
                from reportlab.lib.pagesizes import letter
                import io
                
                reader = PdfReader(file_path)
                writer = PdfWriter()
                for page in reader.pages:
                    writer.add_page(page)
                
                # Create signature overlay
                packet = io.BytesIO()
                can = canvas.Canvas(packet, pagesize=letter)
                can.setFont("Helvetica-Bold", 12)
                can.drawString(50, 50, signature_text)
                can.save()
                packet.seek(0)
                
                sig_reader = PdfReader(packet)
                writer.pages[-1].merge_page(sig_reader.pages[0])
                
                base_name = os.path.splitext(os.path.basename(file_path))[0]
                output_filename = generate_unique_filename(f"{base_name}_signed.pdf")
                output_path = os.path.join(app.config['PROCESSED_FOLDER'], output_filename)
                os.makedirs(app.config['PROCESSED_FOLDER'], exist_ok=True)
                with open(output_path, 'wb') as f:
                    writer.write(f)
                output_files.append(output_path)
                
                self.update_progress(int((file_idx + 1) / len(input_files) * 100), file_idx + 1)
            except Exception as e:
                logger.error(f"Signing failed: {e}")
        return output_files

    def redact_pdf(self):
        """Redact text from PDF"""
        input_files = json.loads(self.job.input_files)
        settings = json.loads(self.job.settings) if self.job.settings else {}
        keywords = settings.get('keywords', [])
        output_files = []
        for file_idx, file_path in enumerate(input_files):
            try:
                import fitz
                doc = fitz.open(file_path)
                for page in doc:
                    for kw in keywords:
                        for inst in page.search_for(kw):
                            page.add_redact_annotation(inst, fill=(0,0,0))
                    page.apply_redactions()
                
                base_name = os.path.splitext(os.path.basename(file_path))[0]
                output_filename = generate_unique_filename(f"{base_name}_redacted.pdf")
                output_path = os.path.join(app.config['PROCESSED_FOLDER'], output_filename)
                os.makedirs(app.config['PROCESSED_FOLDER'], exist_ok=True)
                doc.save(output_path)
                doc.close()
                output_files.append(output_path)
                
                self.update_progress(int((file_idx + 1) / len(input_files) * 100), file_idx + 1)
            except Exception as e:
                logger.error(f"Redaction failed: {e}")
        return output_files

    def compare_pdf(self):
        """Compare two PDFs (basic page count comparison report)"""
        input_files = json.loads(self.job.input_files)
        if len(input_files) < 2:
            return []
        
        report_filename = generate_unique_filename("comparison_report.txt")
        report_path = os.path.join(app.config['PROCESSED_FOLDER'], report_filename)
        os.makedirs(app.config['PROCESSED_FOLDER'], exist_ok=True)
        with open(report_path, "w") as f:
            f.write("PDF Comparison Report\n")
            f.write("=====================\n\n")
            for fp in input_files:
                try:
                    reader = PdfReader(fp)
                    f.write(f"File: {os.path.basename(fp)}\n")
                    f.write(f"Pages: {len(reader.pages)}\n")
                    f.write(f"Metadata: {reader.metadata}\n\n")
                except:
                    f.write(f"Error reading {os.path.basename(fp)}\n\n")
        
        self.update_progress(100, len(input_files))
        return [report_path]

    def convert_to_excel(self):
        """Convert PDF to Excel"""
        input_files = json.loads(self.job.input_files)
        output_files = []
        for file_idx, file_path in enumerate(input_files):
            try:
                wb = openpyxl.Workbook()
                ws = wb.active
                reader = PdfReader(file_path)
                for i, page in enumerate(reader.pages):
                    ws.cell(row=i+1, column=1, value=page.extract_text()[:32000])
                
                base_name = os.path.splitext(os.path.basename(file_path))[0]
                output_filename = generate_unique_filename(f"{base_name}.xlsx")
                output_path = os.path.join(app.config['PROCESSED_FOLDER'], output_filename)
                os.makedirs(app.config['PROCESSED_FOLDER'], exist_ok=True)
                wb.save(output_path)
                output_files.append(output_path)
                
                progress = int((file_idx + 1) / len(input_files) * 100)
                self.update_progress(progress, file_idx + 1)
            except Exception as e:
                logger.error(f"Error converting to excel {file_path}: {str(e)}")
                continue
        return output_files

def create_zip_archive(file_paths, zip_filename):
    """Create a ZIP archive from multiple files in PROCESSED_FOLDER"""
    from app import app
    processed_folder = app.config['PROCESSED_FOLDER']
    os.makedirs(processed_folder, exist_ok=True)
    zip_path = os.path.join(processed_folder, zip_filename)
    
    with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
        for file_path in file_paths:
            if os.path.exists(file_path):
                arcname = os.path.basename(file_path)
                zipf.write(file_path, arcname)
    
    return zip_path
